import io
import os
import subprocess
import tempfile
from pathlib import Path

TEXT_EXTS = {'.txt', '.md', '.csv', '.json', '.log', '.xml', '.html', '.py', '.js', '.ts', '.dart', '.yaml', '.yml', '.cfg', '.ini', '.env', '.sh', '.bat', '.ps1'}
PDF_EXTS = {'.pdf'}
DOCX_EXTS = {'.docx'}
XLSX_EXTS = {'.xlsx', '.xlsm'}
PPTX_EXTS = {'.pptx'}
LIBRE_EXTS = {'.doc', '.xls', '.ppt', '.odt', '.ods', '.odp', '.rtf', '.wpd'}


async def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    try:
        if ext in TEXT_EXTS:
            return _read_text(file_bytes)
        elif ext in PDF_EXTS:
            return _extract_pdf(file_bytes)
        elif ext in DOCX_EXTS:
            return _extract_docx(file_bytes)
        elif ext in XLSX_EXTS:
            return _extract_xlsx(file_bytes)
        elif ext in PPTX_EXTS:
            return _extract_pptx(file_bytes)
        elif ext in LIBRE_EXTS:
            return _extract_libreoffice(file_bytes, filename)
        else:
            return _extract_libreoffice(file_bytes, filename)
    except Exception as e:
        return f"[文件内容提取失败: {e}]"


def _read_text(file_bytes: bytes) -> str:
    for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode('utf-8', errors='replace')


def _extract_pdf(file_bytes: bytes) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            parts.append(text)
        if i >= 50:
            parts.append("...(内容过长，仅提取前50页)")
            break
    return "\n\n".join(parts) if parts else ""


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts[:2000])


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"【Sheet: {name}】")
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                parts.append(row_text)
                row_count += 1
                if row_count >= 500:
                    parts.append("...(该Sheet内容过长，仅提取前500行)")
                    break
    wb.close()
    return "\n".join(parts[:3000])


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"【Slide {i+1}】\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_libreoffice(file_bytes: bytes, filename: str) -> str:
    soffice = os.environ.get("SOFFICE_PATH", "soffice")
    with tempfile.TemporaryDirectory() as tmpdir:
        src = os.path.join(tmpdir, filename)
        with open(src, "wb") as f:
            f.write(file_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, src],
                timeout=60, check=True, capture_output=True,
            )
        except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError):
            return ""
        txt_name = Path(filename).stem + ".txt"
        txt_path = os.path.join(tmpdir, txt_name)
        if not os.path.exists(txt_path):
            return ""
        with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return text[:10000] if len(text) > 10000 else text
